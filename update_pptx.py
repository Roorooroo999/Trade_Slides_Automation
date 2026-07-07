"""
update_pptx.py — Auto-update trade slides PPTX with latest BQ inventory numbers.

Usage:
    python update_pptx.py
    (runs after generate_talk_track_pdf.py pulls fresh BQ data)

How it works:
  1. Pulls the same BQ numbers as the talk track PDF
  2. Builds a find→replace map (old number → new number)
  3. Walks every text run in the PPTX, replaces in-place (preserves fonts/colors)
  4. Saves as  Trade Slides - Inventory WK{TRADE_WK}.pptx  (keeps original safe)

IMPORTANT: Close the PPTX in PowerPoint before running, otherwise PermissionError.
"""

import warnings; warnings.filterwarnings('ignore')
import io, os, sys, re, shutil, datetime
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from dotenv import load_dotenv; load_dotenv()

from google.cloud import bigquery
from pptx import Presentation
from pptx.util import Pt
from business_rules import yyyyww_to_yywww, OO_ROLLING_WEEKS

# ── BQ setup ──────────────────────────────────────────────────────────────────
print("Pulling BQ data for PPTX update...")
client = bigquery.Client(project='wmt-execution-intel-prod')
def bq(sql): return [dict(r) for r in client.query(sql).result()]

# Dynamic week detection (same logic as talk track PDF)
_meta = bq("SELECT MAX(WM_YR_WK_NBR) wk, MAX(BUS_DT) bus_dt "
            "FROM `wmt-execution-intel-prod.WM_AD_HOC.R0C0JUG_WMUS_HIST_COMBINED`")[0]
INV_TY = int(_meta['wk']); INV_LY = INV_TY - 100
OO_TY  = yyyyww_to_yywww(INV_TY)
OO_LY  = OO_TY - 100; OO_PW = OO_TY - 1
DATA_WK  = INV_TY % 100
TRADE_WK = DATA_WK + 1
print(f"  INV_TY={INV_TY}, OO_TY={OO_TY}, TRADE_WK={TRADE_WK}")

# ── Pull node totals ───────────────────────────────────────────────────────────
def _nodes(wk):
    q = f"""WITH s AS (SELECT OMNI_CATG_NBR,MAX(BUS_DT) dt
            FROM `wmt-execution-intel-prod.WM_AD_HOC.R0C0JUG_WMUS_HIST_COMBINED`
            WHERE WM_YR_WK_NBR={wk} GROUP BY 1)
    SELECT CAST(SUM(c.ON_YARD_UNITS) AS FLOAT64) yard,
           CAST(SUM(c.DC_OH_UNITS+c.DC_LABELED_UNITS+c.DC_UNLABELED_UNITS+c.DC_RESERVED_UNITS) AS FLOAT64) dc,
           CAST(SUM(c.INTRANSIT_TO_DC_UNITS) AS FLOAT64) it_dc,
           CAST(SUM(c.IN_TRANSIT_UNITS) AS FLOAT64) it_store,
           CAST(SUM(c.STORE_OH_UNITS) AS FLOAT64) store,
           CAST(SUM(c.BACKROOM_UNITS) AS FLOAT64) backroom,
           CAST(SUM(c.FC_OH_UNITS) AS FLOAT64) fc,
           CAST(SUM(c.TOTAL_NETWORK_UNITS) AS FLOAT64) total_net
    FROM `wmt-execution-intel-prod.WM_AD_HOC.R0C0JUG_WMUS_HIST_COMBINED` c
    INNER JOIN s ON c.OMNI_CATG_NBR=s.OMNI_CATG_NBR AND c.BUS_DT=s.dt"""
    return bq(q)[0]

def _oo(wk):
    return bq(f"SELECT CAST(SUM(units_ordered) AS FLOAT64) v "
              f"FROM `wmt-execution-intel-prod.WM_AD_HOC.R0C0JUG_WMUS_HIST_ONORDER` "
              f"WHERE wm_week={wk} AND dsd_ind='NON-DSD'")[0]['v']

def _l13w(wk, n=OO_ROLLING_WEEKS):
    return bq(f"SELECT CAST(SUM(units_ordered)/{n} AS FLOAT64) v "
              f"FROM `wmt-execution-intel-prod.WM_AD_HOC.R0C0JUG_WMUS_HIST_ONORDER` "
              f"WHERE wm_week BETWEEN {wk-n+1} AND {wk} AND dsd_ind='NON-DSD' AND sbu!='OTHER'")[0]['v']

ty = _nodes(INV_TY); ly = _nodes(INV_LY); pw = _nodes(INV_TY - 1)
oo_ty = _oo(OO_TY); oo_ly = _oo(OO_LY); oo_pw = _oo(OO_PW)
l13w_ty = _l13w(OO_TY); l13w_ly = _l13w(OO_TY - 100)

sf_ty = ty['store'] - ty['backroom']
sf_ly = ly['store'] - ly['backroom']
total_ty = ty['total_net'] + oo_ty
total_ly = ly['total_net'] + oo_ly
pw_total = pw['total_net'] + oo_pw

# ── Format helpers ────────────────────────────────────────────────────────────
def fm(n):
    a = abs(n)
    if a >= 1e9:  return f"{n/1e9:.2f}B"
    if a >= 1e6:  return f"{n/1e6:.1f}M"
    if a >= 1e3:  return f"{n/1e3:.1f}K"
    return str(int(n))

def fp(v): s = "+" if v >= 0 else ""; return f"{s}{v:.1f}%"
def pct(a, b): return ((a-b)/b*100) if b else 0
def dlt(a, b): return a - b

# Slide-format helpers (must be defined before REPLACEMENTS dict)
def _sm(v):
    m = round(v / 1e6)
    return (f"{m:,} M" if m >= 1000 else f"{m} M")

def _sm0(v):
    m = round(v / 1e6)
    return (f"{m:,}M" if m >= 1000 else f"{m}M")

def _sp(v, suffix=" vs LY"):
    if v >= 0: return f"+{v:.1f}%{suffix}"
    return f"({abs(v):.1f}%){suffix}"

def _sp_ly(v): return _sp(v, " vs LY")
def _sp_lw(v): return _sp(v, " vs LW")

# Derived slide values
it_dc_ty = ty['it_dc']; it_dc_ly = ly['it_dc']; it_dc_pw = pw['it_dc']
it_tot_ty = ty['it_dc'] + ty['it_store']
it_tot_ly = ly['it_dc'] + ly['it_store']
it_tot_pw = pw['it_dc'] + pw['it_store']
sf_pw = pw['store'] - pw['backroom']

# ── Build find→replace map ─────────────────────────────────────────────────────
# Format: { "exact text in slide" : "replacement text" }
# Keys must match EXACTLY what appears in the slide text boxes.
# Tip: run with DRY_RUN=True first to print all slide text, then add your mappings.

REPLACEMENTS = {
    # ── Date ───────────────────────────────────────────────────────────────────
    "Trade Meeting · July 6, 2026":     "Trade Meeting · July 13, 2026",

    # ── Summary banner (top row) ───────────────────────────────────────────────
    "TY 8.32 B vs LY 8.54 B":          f"TY {total_ty/1e9:.2f} B vs LY {total_ly/1e9:.2f} B",
    "-2.5% , -0.22 B YoY":             f"{pct(total_ty,total_ly):.1f}% , {dlt(total_ty,total_ly)/1e9:.2f} B YoY",
    "4.65 B store inv":                 f"{ty['store']/1e9:.2f} B store inv",
    "-1.4%, -67.1 M YoY":              f"{pct(ty['store'],ly['store']):.1f}%, {dlt(ty['store'],ly['store'])/1e6:.1f} M YoY",

    # ── On-Order ───────────────────────────────────────────────────────────────
    "1,480M":                           _sm0(oo_ty),                         # units
    "(2.4%) vs LY":                     _sp_ly(pct(oo_ty, oo_ly)),           # YoY
    "(6.0%) vs LW":                     _sp_lw(pct(oo_ty, oo_pw)),           # WoW
    "L13W Avg: 1,370M/wk (+1.6% YoY)": f"L13W Avg: {round(l13w_ty/1e6):,}M/wk ({'+' if pct(l13w_ty,l13w_ly)>=0 else ''}{pct(l13w_ty,l13w_ly):.1f}% YoY)",

    # ── On Yard ────────────────────────────────────────────────────────────────
    "0.3 M":                            f"{ty['yard']/1e6:.1f} M",
    "(94.6%) vs LY":                    _sp_ly(pct(ty['yard'], ly['yard'])),
    "(93.7%) vs LW":                    _sp_lw(pct(ty['yard'], pw['yard'])),

    # ── In DC ──────────────────────────────────────────────────────────────────
    "1,962 M":                          _sm(ty['dc']),
    "(6.4%) vs LY":                     _sp_ly(pct(ty['dc'], ly['dc'])),
    "(6.6%) vs LW":                     _sp_lw(pct(ty['dc'], pw['dc'])),

    # ── IT → DC ────────────────────────────────────────────────────────────────
    "15 M":                             _sm(it_dc_ty),
    "+72.0 % vs LY":                    _sp_ly(pct(it_dc_ty, it_dc_ly)),     # note: space before %
    "+11.2% vs LW":                     _sp_lw(pct(it_dc_ty, it_dc_pw)),

    # ── IT Total ───────────────────────────────────────────────────────────────
    "167 M":                            _sm(it_tot_ty),
    "+15.7% vs LY":                     _sp_ly(pct(it_tot_ty, it_tot_ly)),
    "(0.8%) LW":                        _sp_lw(pct(it_tot_ty, it_tot_pw)),   # note: no "vs" in old text

    # ── IT → Store ─────────────────────────────────────────────────────────────
    "152 M":                            _sm(ty['it_store']),
    "+12.1 % vs LY":                    _sp_ly(pct(ty['it_store'], ly['it_store'])),  # note: space before %
    "(1.8%) vs LW":                     _sp_lw(pct(ty['it_store'], pw['it_store'])),

    # ── Store OH ───────────────────────────────────────────────────────────────
    "4,649 M":                          _sm(ty['store']),
    "(1.4%) vs LY":                     _sp_ly(pct(ty['store'], ly['store'])),   # ⚠ regex prevents cascade
    "+0.1% vs LW":                      _sp_lw(pct(ty['store'], pw['store'])),

    # ── Backroom ───────────────────────────────────────────────────────────────
    "422 M":                            _sm(ty['backroom']),
    "+6.7 % vs LY":                     _sp_ly(pct(ty['backroom'], ly['backroom'])),  # note: space before %
    "+3.9% vs LW":                      _sp_lw(pct(ty['backroom'], pw['backroom'])),

    # ── Salesfloor ─────────────────────────────────────────────────────────────
    "4,227 M":                          _sm(sf_ty),
    "(2.2%) vs LY":                     _sp_ly(pct(sf_ty, sf_ly)),               # ⚠ cascade target — safe with regex
    "(0.2%) vs LW":                     _sp_lw(pct(sf_ty, sf_pw)),

    # ── FC ─────────────────────────────────────────────────────────────────────
    # "63 M" unchanged
    "+5.9% vs LY":                      _sp_ly(pct(ty['fc'], ly['fc'])),
    "(1.6%) vs LW":                     _sp_lw(pct(ty['fc'], pw['fc'])),

    # ── Insights bullets (Rounded Rectangle 10) ─────────────────────────────────
    # Bullet 1: BTX pipeline headline — update to reflect WK23 reality
    "• BTX WK23 pipeline confirmed ready.":
        f"• On-Order turns positive; BTX floor set complete.",
    "DC depleting -6.4% YoY as inventory flows out; In Transit → Store surging +12.1% YoY with all 8 SBUs above last year, landing in stores within 3–5 days. HARDLINES Stationery BTS floor set staging now.":
        (f"On-Order {_sm0(oo_ty)} ({_sp_ly(pct(oo_ty,oo_ly))}) — forward pipeline above LY. "
         f"DC releasing at {_sm(ty['dc'])} ({_sp_ly(pct(ty['dc'],ly['dc']))}). "
         f"IT→Store {_sm(ty['it_store'])} ({_sp_ly(pct(ty['it_store'],ly['it_store']))}) — "
         f"normalizing after last week’s BTS surge; floor set has landed in stores."),

    # Bullet 2: Backroom pull opportunity — update numbers
    "• 422M units in backroom — execution pull opportunity.":
        f"• {_sm(ty['backroom'])} in backroom ({_sp_ly(pct(ty['backroom'],ly['backroom']))}) — sustained pull opportunity.",
    "Store salesfloor is -2.2% YoY (variance in LY data), but backroom inventory is +6.7% YoY. CONSUMABLES, PANTRY, and CAC have product staged and ready to pull to shelf ":
        (f"Salesfloor {_sm(sf_ty)} ({_sp_ly(pct(sf_ty,sf_ly))}) while backroom builds. "
         f"CONSUMABLES, PANTRY, and CAC have product staged — "
         f"pull to shelf without additional receipts. Store OH {_sm(ty['store'])} ({_sp_ly(pct(ty['store'],ly['store']))})."),
}

# ── Print all slide text first so you can fill in REPLACEMENTS above ──────────
DRY_RUN = False   # Set to False once REPLACEMENTS are confirmed

SRC = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   f"Trade Slides - Inventory WK{TRADE_WK}.pptx")
# Falls back to WK23 template if WK{n} doesn't exist yet (first run)
if not os.path.exists(SRC):
    SRC = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       f"Trade Slides - Inventory WK{TRADE_WK-1}.pptx")
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   f"Trade Slides - Inventory WK{TRADE_WK}.pptx")

if not os.path.exists(SRC):
    print(f"ERROR: Source file not found: {SRC}")
    sys.exit(1)

# ── Format helpers matching slide's exact style ───────────────────────────────
def _sm(v):
    """Slide M format: >=1000M uses comma+space+M (e.g. '4,649 M'), else int+space+M."""
    m = round(v / 1e6)
    if m >= 1000: return f"{m:,} M"
    return f"{m} M"

def _sm0(v):
    """On-Order style: no space before M, comma thousands (e.g. '1,480M')."""
    m = round(v / 1e6)
    if m >= 1000: return f"{m:,}M"
    return f"{m}M"

def _sp(v, suffix=" vs LY"):
    """Slide pct: positive = +X.X%, negative = (X.X%)."""
    if v >= 0: return f"+{v:.1f}%{suffix}"
    return f"({abs(v):.1f}%){suffix}"

def _sp_lw(v): return _sp(v, " vs LW")
def _sp_ly(v): return _sp(v, " vs LY")

# Derived values
it_dc_ty = ty['it_dc']; it_dc_ly = ly['it_dc']; it_dc_pw = pw['it_dc']
it_tot_ty = ty['it_dc'] + ty['it_store']
it_tot_ly = ly['it_dc'] + ly['it_store']
it_tot_pw = pw['it_dc'] + pw['it_store']
sf_pw = pw['store'] - pw['backroom']

# Summary banner helper
_tot_delta = dlt(total_ty, total_ly)
_tot_pct   = pct(total_ty, total_ly)
_store_delta = dlt(ty['store'], ly['store'])
_new_banner = (
    f"6 buckets   ·   TY {ty['store']/1e9:.2f} B vs LY {ly['store']/1e9:.2f} B   ·   "
    f"{_tot_pct:.1f}% , {_tot_delta/1e9:.2f} B YoY   ·   "
    f"{ty['store']/1e9:.2f} B store inv ·  "
    f"{pct(ty['store'],ly['store']):.1f}%, {_store_delta/1e6:.1f} M YoY"
)

# Read via open() first to force OneDrive to download cloud-only placeholder files
with open(SRC, "rb") as _f:
    prs = Presentation(io.BytesIO(_f.read()))

if DRY_RUN:
    print("\n" + "="*70)
    print("DRY RUN — all text found in slides:")
    print("="*70)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full = "".join(r.text for r in para.runs)
                    if full.strip():
                        print(f"  [{shape.name}]  {full.strip()[:160]}")
    print("\n" + "="*70)
    print("BQ VALUES (WK{})".format(DATA_WK))
    print("="*70)
    print(f"  On-Order:  {_sm0(oo_ty)} | YoY:{_sp_ly(pct(oo_ty,oo_ly))} | WoW:{_sp_lw(pct(oo_ty,oo_pw))}")
    print(f"  L13W:      1,{round(l13w_ty/1e6):,}M/wk (+{pct(l13w_ty,l13w_ly):.1f}% YoY)")
    print(f"  On Yard:   {ty['yard']/1e6:.1f} M | YoY:{_sp_ly(pct(ty['yard'],ly['yard']))} | WoW:{_sp_lw(pct(ty['yard'],pw['yard']))}")
    print(f"  In DC:     {_sm(ty['dc'])} | YoY:{_sp_ly(pct(ty['dc'],ly['dc']))} | WoW:{_sp_lw(pct(ty['dc'],pw['dc']))}")
    print(f"  IT→DC:     {_sm(it_dc_ty)} | YoY:{_sp_ly(pct(it_dc_ty,it_dc_ly))} | WoW:{_sp_lw(pct(it_dc_ty,it_dc_pw))}")
    print(f"  IT→Store:  {_sm(ty['it_store'])} | YoY:{_sp_ly(pct(ty['it_store'],ly['it_store']))} | WoW:{_sp_lw(pct(ty['it_store'],pw['it_store']))}")
    print(f"  IT Total:  {_sm(it_tot_ty)} | YoY:{_sp_ly(pct(it_tot_ty,it_tot_ly))} | WoW:{_sp_lw(pct(it_tot_ty,it_tot_pw))}")
    print(f"  Store OH:  {_sm(ty['store'])} | YoY:{_sp_ly(pct(ty['store'],ly['store']))} | WoW:{_sp_lw(pct(ty['store'],pw['store']))}")
    print(f"  Backroom:  {_sm(ty['backroom'])} | YoY:{_sp_ly(pct(ty['backroom'],ly['backroom']))} | WoW:{_sp_lw(pct(ty['backroom'],pw['backroom']))}")
    print(f"  Salesfloor:{_sm(sf_ty)} | YoY:{_sp_ly(pct(sf_ty,sf_ly))} | WoW:{_sp_lw(pct(sf_ty,sf_pw))}")
    print(f"  FC:        {_sm(ty['fc'])} | YoY:{_sp_ly(pct(ty['fc'],ly['fc']))} | WoW:{_sp_lw(pct(ty['fc'],pw['fc']))}")
    print(f"  Banner:    {_new_banner}")
    sys.exit(0)

# ── Apply replacements (simultaneous regex — no cascade between keys) ─────────
import re

# Build a single regex that matches all old strings at once
_pattern = re.compile("|".join(re.escape(k) for k in REPLACEMENTS.keys()))

def _replace_all(text):
    """Replace all REPLACEMENTS keys simultaneously (no cascade)."""
    return _pattern.sub(lambda m: REPLACEMENTS[m.group(0)], text)

changes = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            # Strategy: try run-by-run first; also try full-paragraph merge for
            # values that span multiple runs (e.g. banner text box).
            para_text = "".join(r.text for r in para.runs)
            new_para  = _replace_all(para_text)

            if new_para != para_text:
                if len(para.runs) == 1:
                    # Single run — simple replace
                    para.runs[0].text = new_para
                    changes += 1
                    print(f"  [single]  [{shape.name}]  '{para_text[:80]}' → '{new_para[:80]}'")
                else:
                    # Try replacing within individual runs first
                    run_changed = False
                    for run in para.runs:
                        new_run = _replace_all(run.text)
                        if new_run != run.text:
                            run.text = new_run
                            run_changed = True
                    if run_changed:
                        changes += 1
                        print(f"  [per-run] [{shape.name}]  '{para_text[:80]}' → '{''.join(r.text for r in para.runs)[:80]}'")
                    else:
                        # String spans multiple runs (e.g. '(2.4%)' + ' vs LY' in separate runs)
                        # Merge: put full replaced text into first run, clear the rest
                        # Preserves first run's font/color — acceptable since metrics use consistent formatting
                        para.runs[0].text = new_para
                        for run in para.runs[1:]:
                            run.text = ""
                        changes += 1
                        print(f"  [merged]  [{shape.name}]  '{para_text[:80]}' → '{new_para[:80]}'")

prs.save(OUT)
print(f"\n✓  {changes} text run(s) updated.")
print(f"✓  Saved: {OUT}")
