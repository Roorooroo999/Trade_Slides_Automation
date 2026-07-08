"""
Inventory Health Dashboard — Dash app (local + Posit Connect)
"Current inventory units across 6 links of the supply chain"

Run locally:
    python app.py
    # open http://127.0.0.1:8050

Packages required (all already installed):
    dash, dash-bootstrap-components, plotly, pandas,
    google-cloud-bigquery, google-auth, db-dtypes, python-dotenv
"""

from __future__ import annotations

import io, os, re, glob
from datetime import date, datetime, timezone

import dash
import dash_bootstrap_components as dbc
import pandas as pd
import plotly.graph_objects as go
from dash import Input, Output, State, dcc, html
from dotenv import load_dotenv

from data.bq import (
    fetch_inv_snapshot, fetch_oo_snapshot,
    fetch_inv_trend,    fetch_oo_trend,
    fetch_oo_instore_l4w,
    get_depts, get_sbus,
)

load_dotenv()

# ── Colours ───────────────────────────────────────────────────────────────────
WM_BLUE   = "#0071ce"
WM_DARK   = "#003087"
WM_YELLOW = "#ffc220"
WM_GREEN  = "#16a34a"
WM_RED    = "#dc2626"
WM_GREY   = "#f5f5f5"

# ── Data refresh schedule (Wed=2 … Sat=5 in weekday(), Mon=0) ─────────────────
# Data now available daily (scheduled query runs daily Mon-Fri)
_REFRESH_DAYS = {0, 1, 2, 3, 4}       # Mon, Tue, Wed, Thu, Fri
_REFRESH_INTERVAL_MS = 30 * 60 * 1000        # check every 30 minutes

def _should_refresh() -> bool:
    """True Mon-Fri — new pipeline runs daily."""
    return date.today().weekday() in _REFRESH_DAYS

# ── Initial data load ─────────────────────────────────────────────────────────
print("Loading inventory data from BigQuery...")
_INV_DF   = fetch_inv_snapshot()
_OO_DF    = fetch_oo_snapshot()
_CUR_INV_WK = int(_INV_DF["WM_YR_WK_NBR"].max()) if "WM_YR_WK_NBR" in _INV_DF.columns else 202622
_CUR_OO_WK  = int(_OO_DF["wm_week"].max())        if "wm_week"       in _OO_DF.columns  else 12622
_OO_L4W_DF  = fetch_oo_instore_l4w(_CUR_OO_WK, win=2)  # L4W avg + in-store ±2wk
_LAST_LOADED = datetime.now(timezone.utc)
_SBUS  = ["All SBUs"]        + get_sbus(_INV_DF)
_DEPTS = ["All Departments"] + get_depts(_INV_DF)
print(f"Ready. OO week={_CUR_OO_WK} ({len(_OO_DF)} rows), INV week={_CUR_INV_WK} ({len(_INV_DF)} rows)")

# ── Helpers ───────────────────────────────────────────────────────────────────

def _fmt(n: float, _unused: int = 0) -> str:
    """Format numbers with smart unit selection.
    >= 10M  → integer M  e.g. 4,643M / 169M
    >= 1M   → 1-decimal M  e.g. 1.1M / 63.6M
    >= 1K   → K  e.g. 850K / 12K   (for small cube values like On Yard)
    < 1K    → integer
    """
    a = abs(n)
    if a >= 10_000_000:
        return f"{n/1_000_000:,.0f}M"
    if a >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if a >= 1_000:
        return f"{n/1_000:.1f}K"
    return str(int(n))


def _pct(v: float | None, suffix: str = "") -> str:
    if v is None or (isinstance(v, float) and pd.isna(v)):
        return "N/A"
    sign = "+" if v >= 0 else ""
    return f"{sign}{v*100:.1f}%{suffix}"


def _s(df: pd.DataFrame, col: str) -> float:
    return float(df[col].sum()) if col in df.columns else 0.0


def _wpct(df: pd.DataFrame, ty_col: str, prior_col: str) -> float | None:
    """Correct pct: (sum_TY - sum_Prior) / sum_Prior — never average per-row pcts."""
    ty = _s(df, ty_col)
    pw = _s(df, prior_col)
    if pw == 0: return None
    return (ty - pw) / pw


def _wpct_sum2(df: pd.DataFrame, ty_a: str, ty_b: str, ly_a: str, ly_b: str) -> float | None:
    """YoY pct for a difference: (TY_A - TY_B) vs (LY_A - LY_B).
    Used for Salesfloor cube = Store cube - Backroom cube."""
    ty = _s(df, ty_a) - _s(df, ty_b)
    ly = _s(df, ly_a) - _s(df, ly_b)
    if ly == 0: return None
    return (ty - ly) / ly


def _mp(df: pd.DataFrame, col: str) -> float | None:
    """Legacy — only use for columns with no raw prior value available."""
    if col not in df.columns: return None
    v = df[col].dropna()
    return float(v.mean()) if len(v) else None


def _pct_badge(val: float | None, label: str = "") -> html.Span:
    if val is None or (isinstance(val, float) and pd.isna(val)):
        return html.Span("N/A", style={"color": "#aaa", "fontSize": "0.72rem"})
    color = WM_GREEN if val >= 0 else WM_RED
    arrow = "▲" if val >= 0 else "▼"
    sign  = "+" if val >= 0 else ""
    text  = f"{arrow} {sign}{val*100:.1f}%{(' ' + label) if label else ''}"
    return html.Span(text, style={"color": color, "fontWeight": "600", "fontSize": "0.72rem"})


# ── Node card ─────────────────────────────────────────────────────────────────

def node_card(
    num: int, icon: str, title: str, desc: str,
    units: float,
    yoy_pct: float | None, yoy_delta: float | None,
    wow_pct: float | None, wow_delta: float | None,
    cube: float | None = None, cube_yoy: float | None = None,
    l13w_avg: float | None = None, l13w_pct: float | None = None,
    sub_rows: list[tuple] | None = None,
    accent: str = WM_BLUE,
) -> dbc.Card:

    sub_html = []
    if sub_rows:
        # sub_rows tuples: (label, units, wow_pct, yoy_pct)
        #              or: (label, units, wow_pct, yoy_pct, cube, cube_yoy_pct)
        for row in sub_rows:
            lbl, su, sw, sy = row[0], row[1], row[2], row[3]
            sc   = row[4] if len(row) > 4 else None   # cube value
            sc_y = row[5] if len(row) > 5 else None   # cube YoY pct

            cube_line = []
            if sc is not None and sc > 0:
                cube_line = [html.Div([
                    html.Span(f"📦 {_fmt(sc)} ft³", style={"fontSize": "0.68rem", "color": "#aaa"}),
                    html.Span("  "),
                    _pct_badge(sc_y, "vs LY") if sc_y is not None else html.Span(),
                ], style={"marginTop": "1px"})]

            sub_html.append(html.Div([
                html.Div([
                    html.Span(lbl, style={"fontWeight": "600", "fontSize": "0.73rem", "color": "#444"}),
                    html.Span(_fmt(su), style={"float": "right", "fontWeight": "700", "fontSize": "0.73rem"}),
                ], style={"overflow": "hidden"}),
                html.Div([
                    _pct_badge(sy, "YoY") if sy is not None else html.Span(),
                    html.Span("  ") if sy is not None and sw is not None else html.Span(),
                    _pct_badge(sw, "WoW") if sw is not None else html.Span(),
                ]),
                *cube_line,
            ], style={
                "background": "#f8fafc", "borderRadius": "4px",
                "padding": "4px 6px", "marginTop": "4px",
            }))

    cube_row = []
    if cube is not None:
        cube_row = [html.Div([
            html.Span(f"📦 {_fmt(cube)} ft³", style={"fontSize": "0.72rem", "color": "#888"}),
            html.Span("  "),
            _pct_badge(cube_yoy, "vs LY"),
        ], style={"borderTop": "1px solid #eee", "paddingTop": "4px", "marginTop": "6px"})]

    l13w_row = []
    if l13w_avg is not None:
        l13w_row = [html.Div([
            html.Span(f"L13W Avg: ", style={"color": "#555", "fontSize": "0.71rem"}),
            html.Strong(_fmt(l13w_avg) + "/wk", style={"fontSize": "0.71rem"}),
            html.Span("  "),
            _pct_badge(l13w_pct, "vs LY"),
        ], style={
            "background": "#eef4ff", "borderRadius": "4px",
            "padding": "3px 6px", "marginTop": "4px", "fontSize": "0.71rem",
        })]

    return dbc.Card(
        dbc.CardBody([
            # Header
            html.Div([
                html.Span(str(num), style={
                    "background": accent, "color": "white",
                    "borderRadius": "50%", "width": "20px", "height": "20px",
                    "display": "inline-flex", "alignItems": "center",
                    "justifyContent": "center", "fontSize": "0.62rem",
                    "fontWeight": "800", "flexShrink": "0",
                }),
                html.Span(icon, style={"fontSize": "0.65rem", "color": "#bbb", "marginLeft": "4px"}),
                html.Span(title, style={"fontWeight": "700", "fontSize": "0.72rem", "marginLeft": "4px"}),
            ], style={"display": "flex", "alignItems": "center", "marginBottom": "4px"}),
            html.P(desc, style={"fontSize": "0.64rem", "color": "#999", "marginBottom": "6px", "lineHeight": "1.3"}),
            # Main value
            html.Div(_fmt(units), style={"fontSize": "1.5rem", "fontWeight": "800", "color": accent, "lineHeight": "1.1"}),
            html.Div([_pct_badge(yoy_pct), html.Span(" vs LY", style={"fontSize": "0.68rem", "color": "#777"})]),
            html.Div([_pct_badge(wow_pct), html.Span(" vs LW", style={"fontSize": "0.68rem", "color": "#777"})]),
            *l13w_row,
            *sub_html,
            *cube_row,
        ], style={"padding": "12px 14px"}),
        style={
            "borderTop": f"3px solid {accent}",
            "borderRadius": "8px",
            "boxShadow": "0 1px 4px rgba(0,0,0,.10)",
            "height": "100%",
        },
    )


# ── App ───────────────────────────────────────────────────────────────────────
app = dash.Dash(
    __name__,
    external_stylesheets=[dbc.themes.BOOTSTRAP],
    title="Inventory Health Dashboard",
)
server = app.server   # for Posit Connect / gunicorn

# ── Layout ────────────────────────────────────────────────────────────────────
app.layout = dbc.Container(
    fluid=True,
    style={"maxWidth": "1480px", "padding": "20px 24px", "backgroundColor": WM_GREY},
    children=[

        # ── Hidden stores + auto-refresh interval ───────────────────────────
        dcc.Store(id="insights-store", data=""),
        dcc.Store(id="data-store", data={"loaded": str(_LAST_LOADED)}),
        # Fires every 4 hours — callback checks if today is Wed–Sat before reloading
        dcc.Interval(
            id="refresh-interval",
            interval=_REFRESH_INTERVAL_MS,
            n_intervals=0,
            disabled=not _should_refresh(),
        ),

        # ── Header row ───────────────────────────────────────────────────────
        dbc.Row([
            dbc.Col([
                html.Div("Merch Strategy & Enablement | Inventory Insights",
                         style={"fontSize": "0.68rem", "color": "#888", "fontWeight": "500"}),
                html.H5([
                    html.Strong("Current inventory units "),
                    html.Span("across 6 links of the supply chain", style={"fontWeight": "400"}),
                ], style={"margin": "2px 0 4px 0", "fontSize": "1.2rem"}),
                html.Div([
                    html.Span("INVENTORY UNITS AND CUBIC ON HAND · CURRENT STATE",
                              style={"marginRight": "10px"}),
                    html.Span(id="last-updated-lbl",
                              style={"fontSize": "0.65rem", "color": "#888", "fontStyle": "italic"}),
                ], style={
                    "display": "inline-flex", "alignItems": "center",
                    "background": "#e8f0fe", "color": WM_BLUE, "borderRadius": "4px",
                    "padding": "3px 10px", "fontSize": "0.68rem",
                    "fontWeight": "700", "letterSpacing": "0.04em",
                }),
            ], width=8, style={"display": "flex", "flexDirection": "column", "justifyContent": "flex-end"}),
            dbc.Col([
                dbc.Row([
                    dbc.Col(dcc.Dropdown(
                        id="sbu-filter",
                        options=[{"label": s, "value": s} for s in _SBUS],
                        value="All SBUs", clearable=False,
                        style={"fontSize": "0.82rem"},
                    ), width=6),
                    dbc.Col(dcc.Dropdown(
                        id="dept-filter",
                        options=[{"label": d, "value": d} for d in _DEPTS],
                        value="All Departments", clearable=False,
                        style={"fontSize": "0.82rem"},
                    ), width=6),
                ], className="g-2"),
                dbc.Row([
                    dbc.Col([
                        dbc.Button(
                            "⬇ Download Trade Slides (PPTX)",
                            id="btn-download-pptx", n_clicks=0,
                            style={
                                "background": WM_DARK, "border": "none",
                                "fontWeight": "700", "color": "white",
                                "fontSize": "0.78rem", "width": "100%",
                                "marginTop": "6px",
                            },
                        ),
                        dcc.Download(id="download-pptx"),
                    ], width=12),
                ], className="g-2"),
            ], width=4, className="d-flex align-items-end flex-column"),
        ], className="mb-3"),

        # ── Summary banner ───────────────────────────────────────────────────
        html.Div(id="summary-banner", style={
            "background": WM_BLUE, "color": "white", "borderRadius": "8px",
            "padding": "12px 20px", "marginBottom": "16px",
            "display": "flex", "alignItems": "center",
            "flexWrap": "wrap", "gap": "20px",
        }),

        # ── 6 Node cards ─────────────────────────────────────────────────────
        dbc.Row([
            dbc.Col(html.Div(id="card-factory"),  width=2),
            dbc.Col(html.Div(id="card-yard"),      width=2),
            dbc.Col(html.Div(id="card-dc"),        width=2),
            dbc.Col(html.Div(id="card-intransit"), width=2),
            dbc.Col(html.Div(id="card-store"),     width=2),
            dbc.Col(html.Div(id="card-fc"),        width=2),
        ], className="g-2 mb-3"),

        # ── Detail table ─────────────────────────────────────────────────────
        # ── On-Order: MABD vs L4W vs In-Store comparison ─────────────────────
        dbc.Card([
            dbc.CardBody([
                html.Div([
                    html.Strong("🏭 On-Order: MABD vs L4W Rolling Avg vs In-Store Date",
                                style={"fontSize": "0.95rem"}),
                ], className="mb-1"),
                html.P(
                    "MABD single-week view is noisy (import PO timing shifts all land in one week). "
                    "L4W rolling avg smooths this. In-store date shows actual shelf arrival — "
                    "only reliable for SBUs with ≥15% coverage (HOME, ETS, FASHION, HARDLINES).",
                    style={"fontSize": "0.74rem", "color": "#999", "marginBottom": "10px"}),
                html.Div(id="oo-l4w-table"),
            ])
        ], style={"borderRadius": "8px", "boxShadow": "0 1px 4px rgba(0,0,0,.10)",
                  "marginBottom": "16px"}),

        dbc.Card([
            dbc.CardBody([
                html.Div([
                    html.Strong("Inventory Detail by SBU / Department", style={"fontSize": "0.95rem"}),
                    html.Span(" · WoW Δ · YoY Δ · Cube",
                              style={"fontSize": "0.74rem", "color": "#999", "marginLeft": "6px"}),
                ], className="mb-2"),
                dbc.Row([
                    dbc.Col(dcc.Dropdown(
                        id="table-node",
                        options=[
                            {"label": "Store OH",            "value": "store"},
                            {"label": "On-Order",            "value": "on_order"},
                            {"label": "In DC",               "value": "in_dc"},
                            {"label": "In Transit (Total)",  "value": "it_total"},
                            {"label": "In Transit → DC",     "value": "it_dc"},
                            {"label": "In Transit → Store",  "value": "it_store"},
                            {"label": "Backroom",            "value": "backroom"},
                            {"label": "FC",                  "value": "fc"},
                            {"label": "On Yard",             "value": "yard"},
                            {"label": "Total Network",       "value": "total"},
                        ],
                        value="store", clearable=False,
                        style={"fontSize": "0.82rem", "width": "220px"},
                    ), width="auto"),
                ], className="mb-3"),
                html.Div(id="detail-table"),
            ])
        ], style={"borderRadius": "8px", "boxShadow": "0 1px 4px rgba(0,0,0,.10)", "marginBottom": "16px"}),

        # ── Trend chart ──────────────────────────────────────────────────────
        dbc.Card([
            dbc.CardBody([
                html.Strong("Weekly Trend", style={"fontSize": "0.95rem"}),
                dbc.Row([
                    dbc.Col([
                        html.Div("Node:", style={"fontSize": "0.78rem", "color": "#666", "marginBottom": "4px"}),
                        dbc.RadioItems(
                            id="trend-node",
                            options=[
                                {"label": "Store OH",     "value": "store_oh_units"},
                                {"label": "On Yard",      "value": "on_yard_units"},
                                {"label": "In DC",        "value": "in_dc_units"},
                                {"label": "In Transit",   "value": "it_total_units"},
                                {"label": "FC",           "value": "fc_oh_units"},
                                {"label": "Total Net",    "value": "total_network_units"},
                                {"label": "On-Order",     "value": "on_order"},
                            ],
                            value="store_oh_units",
                            inline=True,
                            inputStyle={"marginRight": "4px"},
                            labelStyle={"fontSize": "0.78rem", "marginRight": "14px"},
                        ),
                    ], width=9),
                    dbc.Col([
                        html.Div("Weeks:", style={"fontSize": "0.78rem", "color": "#666", "marginBottom": "4px"}),
                        dbc.RadioItems(
                            id="trend-range",
                            options=[
                                {"label": "8 wks",  "value": "8"},
                                {"label": "13 wks", "value": "13"},
                                {"label": "26 wks", "value": "26"},
                            ],
                            value="13",
                            inline=True,
                            inputStyle={"marginRight": "4px"},
                            labelStyle={"fontSize": "0.78rem", "marginRight": "10px"},
                        ),
                    ], width=3),
                ], className="mt-2 mb-2"),
                dcc.Graph(id="trend-chart", config={"displayModeBar": False}),
            ])
        ], style={"borderRadius": "8px", "boxShadow": "0 1px 4px rgba(0,0,0,.10)", "marginBottom": "16px"}),

        # ── AI Insights ──────────────────────────────────────────────────────
        dbc.Card([
            dbc.CardBody([
                html.Div([
                    html.Span("AI", style={
                        "background": WM_BLUE, "color": "white", "borderRadius": "4px",
                        "padding": "2px 8px", "fontSize": "0.7rem", "fontWeight": "700",
                        "marginRight": "8px",
                    }),
                    html.Strong("Weekly Inventory Insights Agent", style={"fontSize": "0.95rem"}),
                ], className="mb-1"),
                html.P(
                    "Generate an executive-ready weekly update comparing the most recent week "
                    "to the prior week, with cross-node supply chain narrative. Formatted for email.",
                    style={"fontSize": "0.73rem", "color": "#999", "marginBottom": "14px"},
                ),
                dbc.Row([
                    dbc.Col([
                        html.Small("Current WM Week", style={"fontWeight": "600"}),
                        html.Div(id="cur-week-lbl",
                                 style={"fontWeight": "700", "color": WM_BLUE, "fontSize": "0.9rem"}),
                    ], width="auto"),
                    dbc.Col([
                        html.Small("Prior Week", style={"fontWeight": "600"}),
                        html.Div(id="prior-week-lbl",
                                 style={"fontWeight": "700", "color": "#666", "fontSize": "0.9rem"}),
                    ], width="auto"),
                    dbc.Col(
                        dbc.Button("Generate Insights", id="btn-generate", n_clicks=0,
                                   style={"background": WM_YELLOW, "border": "none",
                                          "fontWeight": "700", "color": "#222", "fontSize": "0.83rem"}),
                        width="auto",
                    ),
                    dbc.Col(
                        dbc.Button("Copy for Email", id="btn-copy", n_clicks=0,
                                   outline=True, color="secondary", size="sm",
                                   style={"fontSize": "0.83rem"}),
                        width="auto",
                    ),
                ], align="center", className="g-2 mb-3"),
                html.Div(id="insights-output", children=[
                    html.Em("Click 'Generate Insights' to produce the weekly supply chain executive summary.",
                            style={"color": "#bbb", "fontSize": "0.85rem"}),
                ], style={
                    "background": "#f8fafc", "borderRadius": "6px", "padding": "16px",
                    "border": "1px solid #e2e8f0", "minHeight": "90px",
                    "fontSize": "0.87rem", "lineHeight": "1.75", "whiteSpace": "pre-wrap",
                }),
            ])
        ], style={
            "borderRadius": "8px",
            "border": f"2px solid {WM_BLUE}",
            "boxShadow": "0 1px 4px rgba(0,0,0,.08)",
        }),

        # clipboard helper
        dcc.Clipboard(id="clipboard", style={"display": "none"}),
    ],
)


# ── Callbacks ─────────────────────────────────────────────────────────────────

# Module-level mutable cache (updated by refresh callback)
_cache: dict = {
    "inv":     _INV_DF,
    "oo":      _OO_DF,
    "oo_l4w":  _OO_L4W_DF,
    "cur_inv_wk": _CUR_INV_WK,
    "cur_oo_wk":  _CUR_OO_WK,
    "last_loaded": _LAST_LOADED,
}


def _get_frames(sbu: str, dept: str):
    inv = _cache["inv"].copy()
    oo  = _cache["oo"].copy()
    if sbu  != "All SBUs":        inv = inv[inv["sbu"] == sbu];           oo = oo[oo["sbu"] == sbu]
    if dept != "All Departments": inv = inv[inv["OMNI_DEPT_DESC"] == dept]; oo = oo[oo["OMNI_DEPT_DESC"] == dept]
    return inv, oo


@app.callback(
    Output("oo-l4w-table", "children"),
    Input("data-store", "data"),
)
def render_oo_l4w_table(_):
    df = _cache.get("oo_l4w", pd.DataFrame())
    if df.empty:
        return html.P("In-Store / L4W data not available.", style={"color": "#aaa", "fontStyle": "italic", "fontSize": "0.82rem"})

    def _pc(v, threshold=None):
        """Color-coded pct badge."""
        if v is None or (isinstance(v, float) and pd.isna(v)): return html.Td("—")
        color = WM_GREEN if v >= 0 else WM_RED
        s = "+" if v >= 0 else ""
        style = {"color": color, "fontWeight": "600", "textAlign": "right"}
        return html.Td(f"{s}{v*100:.1f}%", style=style)

    def _td(v, bold=False):
        return html.Td(_fmt(v), style={"textAlign": "right", "fontWeight": "700" if bold else "400"})

    def _cov_badge(pct):
        color = WM_GREEN if pct >= 30 else "#f59e0b" if pct >= 10 else WM_RED
        return html.Td(f"{pct:.0f}%", style={"textAlign": "center", "color": color, "fontWeight": "600", "fontSize": "0.75rem"})

    TH = {"background": "#0071ce", "color": "white", "padding": "7px 10px",
          "fontSize": "0.78rem", "fontWeight": "700", "textAlign": "right",
          "whiteSpace": "nowrap", "borderBottom": "2px solid #003087"}
    TH_L = {**TH, "textAlign": "left"}
    TD_SEP = {"borderLeft": "2px solid #e0e0e0"}

    headers = html.Tr([
        html.Th("SBU",              style=TH_L),
        html.Th("MABD TY",          style=TH),
        html.Th("MABD WoW%",        style=TH),
        html.Th("MABD YoY%",        style=TH),
        html.Th("L4W Avg TY",       style={**TH, **TD_SEP}),
        html.Th("L4W Avg LY",       style=TH),
        html.Th("L4W YoY%",         style=TH),
        html.Th("InStr Cov%",       style={**TH, **TD_SEP}),
        html.Th("InStr ±2wk TY",    style=TH),
        html.Th("InStr ±2wk LY",    style=TH),
        html.Th("InStr YoY%",       style=TH),
    ])

    rows = []
    for _, r in df.iterrows():
        mabd_ty  = float(r.get("mabd_ty",  0) or 0)
        mabd_pw  = float(r.get("mabd_pw",  0) or 0)
        mabd_ly  = float(r.get("mabd_ly",  0) or 0)
        l4w_ty   = float(r.get("l4w_avg_ty", 0) or 0)
        l4w_ly   = float(r.get("l4w_avg_ly", 0) or 0)
        cov      = float(r.get("ins_cov_pct", 0) or 0)
        ins_ty   = float(r.get("instore_win_ty", 0) or 0)
        ins_ly   = float(r.get("instore_win_ly", 0) or 0)
        sbu      = r.get("sbu", "")

        wow  = (mabd_ty - mabd_pw) / mabd_pw if mabd_pw else None
        yoy  = (mabd_ty - mabd_ly) / mabd_ly if mabd_ly else None
        l4y  = (l4w_ty  - l4w_ly)  / l4w_ly  if l4w_ly  else None
        iyoy = (ins_ty  - ins_ly)   / ins_ly   if ins_ly  else None

        rows.append(html.Tr([
            html.Td(sbu, style={"fontWeight": "700", "color": WM_DARK, "padding": "6px 10px", "fontSize": "0.82rem"}),
            _td(mabd_ty, bold=True),
            _pc(wow),
            _pc(yoy),
            html.Td(_fmt(l4w_ty), style={"textAlign": "right", "fontWeight": "700", **TD_SEP}),
            _td(l4w_ly),
            _pc(l4y),
            _cov_badge(cov),
            _td(ins_ty) if cov >= 10 else html.Td("⚠ low cov", style={"textAlign": "center", "color": "#aaa", "fontSize": "0.75rem"}),
            _td(ins_ly) if cov >= 10 else html.Td("—", style={"textAlign": "center", "color": "#aaa"}),
            _pc(iyoy)   if cov >= 10 else html.Td("N/A", style={"textAlign": "center", "color": "#aaa", "fontSize": "0.75rem"}),
        ], style={"borderBottom": "1px solid #f2f2f2", "fontSize": "0.82rem"}))

    return html.Div(
        html.Table(
            [html.Thead(headers), html.Tbody(rows)],
            style={"width": "100%", "borderCollapse": "collapse"}
        ),
        style={"overflowX": "auto"}
    )


# Auto-refresh: reload BQ data on Wed–Sat every 4 hours
# Only updates data-store; update_cards re-fires automatically and sets last-updated-lbl
@app.callback(
    Output("data-store", "data"),
    Input("refresh-interval", "n_intervals"),
    prevent_initial_call=True,
)
def auto_refresh(n):
    if not _should_refresh():
        return dash.no_update

    # ── Smart refresh: check if BQ data actually changed before reloading ────
    # Runs a cheap MAX(BUS_DT) query — avoids reloading 1500+ rows when
    # the scheduled query hasn't landed yet.
    try:
        from data.bq import _run_query
        chk = _run_query(
            "SELECT MAX(BUS_DT) latest FROM "
            "`wmt-execution-intel-prod.WM_AD_HOC.R0C0JUG_WMUS_HIST_COMBINED`"
        )
        if not chk.empty:
            latest_bus_dt = str(chk["latest"].iloc[0])
            last_seen     = _cache.get("last_bus_dt", "")
            if latest_bus_dt == last_seen:
                print(f"[Refresh] BQ unchanged (BUS_DT={latest_bus_dt}) — skipping reload")
                return dash.no_update
            print(f"[Refresh] New BQ data detected: {last_seen} → {latest_bus_dt}")
            _cache["last_bus_dt"] = latest_bus_dt
    except Exception as e:
        print(f"[Refresh] BUS_DT check failed ({e}) — proceeding with full reload")

    print(f"[Refresh] Loading fresh BQ snapshot (interval #{n})...")
    new_inv = fetch_inv_snapshot()
    new_oo  = fetch_oo_snapshot()

    if not new_inv.empty and not new_oo.empty:
        _cache["inv"] = new_inv
        _cache["oo"]  = new_oo
        _cache["cur_inv_wk"] = int(new_inv["WM_YR_WK_NBR"].max()) if "WM_YR_WK_NBR" in new_inv.columns else _cache["cur_inv_wk"]
        _cache["cur_oo_wk"]  = int(new_oo["wm_week"].max())        if "wm_week"       in new_oo.columns  else _cache["cur_oo_wk"]
        new_l4w = fetch_oo_instore_l4w(_cache["cur_oo_wk"])
        if not new_l4w.empty:
            _cache["oo_l4w"] = new_l4w
        _cache["last_loaded"] = datetime.now(timezone.utc)
        print(f"[Refresh] Done. OO wk={_cache['cur_oo_wk']}, INV wk={_cache['cur_inv_wk']}")

        # ── Auto-regenerate Build/Burn chart PNG alongside BQ refresh ─────────
        # Runs in a background thread so it doesn't block the dashboard response
        import threading
        def _regen_chart():
            try:
                import capture_buildburn
                chart_png = os.path.join(_BASE_DIR, "buildburn_chart.png")
                ok = capture_buildburn.capture(out_png=chart_png)
                if ok:
                    print(f"[Refresh] Build/Burn chart auto-updated ({os.path.getsize(chart_png):,} bytes)", flush=True)
                else:
                    print("[Refresh] Build/Burn chart auto-update failed (non-blocking)", flush=True)
            except Exception as e:
                print(f"[Refresh] Build/Burn chart error (non-blocking): {e}", flush=True)
        threading.Thread(target=_regen_chart, daemon=True).start()

    return {"loaded": str(_cache["last_loaded"])}


@app.callback(
    Output("summary-banner",  "children"),
    Output("card-factory",    "children"),
    Output("card-yard",       "children"),
    Output("card-dc",         "children"),
    Output("card-intransit",  "children"),
    Output("card-store",      "children"),
    Output("card-fc",         "children"),
    Output("cur-week-lbl",    "children"),
    Output("prior-week-lbl",  "children"),
    Output("last-updated-lbl","children"),
    Input("sbu-filter",  "value"),
    Input("dept-filter", "value"),
    Input("data-store",  "data"),   # re-fires when data refreshes
)
def update_cards(sbu, dept, _data):
    inv, oo = _get_frames(sbu, dept)

    # ── Summary banner ────────────────────────────────────────────────────
    # Total network = HIST_COMBINED nodes + On-Order (6-bucket total)
    inv_net     = _s(inv, "total_network_units")
    ly_inv_net  = _s(inv, "ly_total")
    oo_u        = _s(oo,  "units_ordered")
    ly_oo_u     = _s(oo,  "ly_units_ordered")

    total_6bkt  = inv_net  + oo_u        # TY: all 6 buckets
    ly_6bkt     = ly_inv_net + ly_oo_u   # LY: all 6 buckets

    yoy_6bkt_pct   = (total_6bkt - ly_6bkt) / ly_6bkt if ly_6bkt else None
    yoy_6bkt_delta = total_6bkt - ly_6bkt               # raw unit delta

    wow_t       = _wpct(inv, "total_network_units", "pw_total")

    # Store
    store       = _s(inv, "store_oh_units")
    ly_store    = _s(inv, "ly_store")
    s_yoy_pct   = (store - ly_store) / ly_store if ly_store else None
    s_yoy_delta = store - ly_store                       # raw unit delta

    oo_yoy      = _wpct(oo, "units_ordered", "ly_units_ordered")

    def _bkpi(lbl, val, sub):
        return html.Div([
            html.Div(val, style={"fontSize": "1.25rem", "fontWeight": "800"}),
            html.Div(lbl, style={"fontSize": "0.65rem", "opacity": ".85",
                                  "textTransform": "uppercase", "letterSpacing": ".04em"}),
            html.Div(sub, style={"fontSize": "0.68rem", "opacity": ".9"}),
        ])

    def _divider():
        return html.Div(style={"width": "1px", "background": "rgba(255,255,255,.3)", "height": "40px"})

    def _delta_str(delta: float) -> str:
        """Format unit delta: +0.03B / -30.2M"""
        sign = "+" if delta >= 0 else ""
        a = abs(delta)
        if a >= 1_000_000_000:
            return f"{sign}{delta/1_000_000_000:.2f}B"
        if a >= 1_000_000:
            return f"{sign}{delta/1_000_000:.1f}M"
        return f"{sign}{delta/1_000:.1f}K"

    banner = [
        html.Div("6 buckets", style={"fontSize": "0.7rem", "fontWeight": "700",
                                      "letterSpacing": ".06em", "textTransform": "uppercase",
                                      "opacity": ".8", "marginRight": "8px"}),
        _bkpi(
            f"TY {total_6bkt/1e9:.2f}B vs LY {ly_6bkt/1e9:.2f}B",
            f"{_pct(yoy_6bkt_pct)}, {_delta_str(yoy_6bkt_delta)} YoY",
            f"{_pct(wow_t)} WoW · incl. On-Order",
        ),
        _divider(),
        _bkpi(
            f"{store/1e9:.2f}B store inv",
            f"{_pct(s_yoy_pct)}, {_delta_str(s_yoy_delta)} YoY",
            "",
        ),
        _divider(),
        _bkpi("On Order",  _fmt(oo_u), f"{_pct(oo_yoy)} YoY"),
        _divider(),
        _bkpi("Current WM Week", f"WK {_cache['cur_oo_wk']}", ""),
    ]

    # ── Node cards (all pcts from raw sums) ───────────────────────────────
    factory = node_card(
        1, "🏭", "FACTORY", "All On-Order from Vendor → DC Gate (WK21 MABDs, NON-DSD)",
        units    = _s(oo, "units_ordered"),
        yoy_pct  = _wpct(oo, "units_ordered",  "ly_units_ordered"),
        yoy_delta= _s(oo, "units_ordered") - _s(oo, "ly_units_ordered"),
        wow_pct  = _wpct(oo, "units_ordered",  "pw_units_ordered"),
        wow_delta= _s(oo, "units_ordered") - _s(oo, "pw_units_ordered"),
        cube     = _s(oo, "cube_ordered"),
        cube_yoy = _wpct(oo, "cube_ordered", "ly_cube_ordered"),
        l13w_avg = _s(oo, "l13w_avg_units"),
        l13w_pct = _wpct(oo, "l13w_avg_units", "l13w_avg_units_ly"),
        accent   = WM_DARK,
    )
    yard = node_card(
        2, "🚚", "YARD – RDC & FDC", "At the Gate — not yet received",
        units    = _s(inv, "on_yard_units"),
        yoy_pct  = _wpct(inv, "on_yard_units", "ly_yard"),
        yoy_delta= _s(inv, "on_yard_units") - _s(inv, "ly_yard"),
        wow_pct  = _wpct(inv, "on_yard_units", "pw_yard"),
        wow_delta= _s(inv, "on_yard_units") - _s(inv, "pw_yard"),
        cube     = _s(inv, "on_yard_cube"),
        cube_yoy = _wpct(inv, "on_yard_cube",  "ly_yard_cube"),
        accent   = "#004f9f",
    )
    dc = node_card(
        3, "🏪", "WAREHOUSE", "DC OH + Labeled + Unlabeled + Reserved",
        units    = _s(inv, "in_dc_units"),
        yoy_pct  = _wpct(inv, "in_dc_units",   "ly_dc"),
        yoy_delta= _s(inv, "in_dc_units") - _s(inv, "ly_dc"),
        wow_pct  = _wpct(inv, "in_dc_units",   "pw_dc"),
        wow_delta= _s(inv, "in_dc_units") - _s(inv, "pw_dc"),
        cube     = _s(inv, "in_dc_cube"),
        cube_yoy = _wpct(inv, "in_dc_cube",    "ly_dc_cube"),
        accent   = WM_BLUE,
    )
    it_total     = _s(inv, "it_total_units")
    pw_it        = _s(inv, "pw_it_dc")      + _s(inv, "pw_it_store")
    ly_it        = _s(inv, "ly_it_dc")      + _s(inv, "ly_it_store")
    it_cube      = _s(inv, "it_total_cube")
    pw_it_cube   = _s(inv, "pw_it_dc_cube") + _s(inv, "pw_it_store_cube")
    ly_it_cube   = _s(inv, "ly_it_dc_cube") + _s(inv, "ly_it_store_cube")
    it = node_card(
        4, "🚛", "IN TRANSIT", "On-truck inbound to DC or outbound to Store",
        units    = it_total,
        yoy_pct  = (it_total - ly_it)   / ly_it   if ly_it   else None,
        yoy_delta= it_total - ly_it,
        wow_pct  = (it_total - pw_it)   / pw_it   if pw_it   else None,
        wow_delta= it_total - pw_it,
        cube     = it_cube,
        cube_yoy = (it_cube - ly_it_cube) / ly_it_cube if ly_it_cube else None,
        sub_rows = [
            ("→ RDC/FDC/ICC/ACC", _s(inv, "it_dc_units"),
             _wpct(inv, "it_dc_units",    "pw_it_dc"),
             _wpct(inv, "it_dc_units",    "ly_it_dc"),
             _s(inv, "it_dc_cube"),
             _wpct(inv, "it_dc_cube",     "ly_it_dc_cube")),
            ("→ Store",           _s(inv, "it_store_units"),
             _wpct(inv, "it_store_units", "pw_it_store"),
             _wpct(inv, "it_store_units", "ly_it_store"),
             _s(inv, "it_store_cube"),
             _wpct(inv, "it_store_cube",  "ly_it_store_cube")),
        ],
        accent = "#4a90d9",
    )
    # Backroom: show full total including Fashion
    # YoY shown in card but note in talk track that Fashion inflates YoY (new FY26 tracking)
    # WoW is valid for all SBUs including Fashion (Fashion WAS in backroom last week)
    br_total_ty  = _s(inv, "backroom_units")
    br_total_ly  = _s(inv, "ly_backroom")
    br_total_wow = _wpct(inv, "backroom_units", "pw_backroom")

    store_card = node_card(
        5, "🛒", "STORE", "Store floor + backroom",
        units    = _s(inv, "store_oh_units"),
        yoy_pct  = _wpct(inv, "store_oh_units", "ly_store"),
        yoy_delta= _s(inv, "store_oh_units") - _s(inv, "ly_store"),
        wow_pct  = _wpct(inv, "store_oh_units", "pw_store"),
        wow_delta= _s(inv, "store_oh_units") - _s(inv, "pw_store"),
        cube     = _s(inv, "store_oh_cube"),
        cube_yoy = _wpct(inv, "store_oh_cube",  "ly_store_cube"),
        sub_rows = [
            ("Backroom",  br_total_ty,
             br_total_wow,
             (br_total_ty - br_total_ly) / br_total_ly if br_total_ly else None,
             _s(inv, "backroom_cube"),
             _wpct(inv, "backroom_cube", "ly_backroom_cube")),
            # Salesfloor = Store OH - Backroom (ensures reconciliation holds)
            # Salesfloor = Store OH − Backroom (consistent identity)
            # YoY and WoW computed the same way: (sf_ty − sf_prior) / sf_prior
            ("Salesfloor",
             _s(inv, "store_oh_units") - _s(inv, "backroom_units"),
             # WoW: (sf_ty - sf_pw) / sf_pw
             _wpct_sum2(inv, "store_oh_units", "backroom_units", "pw_store", "pw_backroom"),
             # YoY: (sf_ty - sf_ly) / sf_ly
             _wpct_sum2(inv, "store_oh_units", "backroom_units", "ly_store", "ly_backroom"),
             _s(inv, "store_oh_cube") - _s(inv, "backroom_cube"),
             _wpct_sum2(inv, "store_oh_cube", "backroom_cube", "ly_store_cube", "ly_backroom_cube")),
        ],
        accent = WM_GREEN,
    )
    fc = node_card(
        6, "📦", "FC", "In FC — online fulfillment",
        units    = _s(inv, "fc_oh_units"),
        yoy_pct  = _wpct(inv, "fc_oh_units", "ly_fc"),
        yoy_delta= _s(inv, "fc_oh_units") - _s(inv, "ly_fc"),
        wow_pct  = _wpct(inv, "fc_oh_units", "pw_fc"),
        wow_delta= _s(inv, "fc_oh_units") - _s(inv, "pw_fc"),
        cube     = _s(inv, "fc_oh_cube"),
        cube_yoy = _wpct(inv, "fc_oh_cube",  "ly_fc_cube"),
        accent   = WM_YELLOW,
    )

    cw = _cache["cur_oo_wk"]
    ts = _cache["last_loaded"].strftime("Updated %b %d %H:%M UTC")
    return (banner, factory, yard, dc, it, store_card, fc,
            f"WK {cw}", f"WK {cw - 1}", ts)


@app.callback(
    Output("detail-table", "children"),
    Input("table-node",  "value"),
    Input("sbu-filter",  "value"),
    Input("dept-filter", "value"),
)
def update_table(node, sbu, dept):
    inv, oo = _get_frames(sbu, dept)

    # Each entry: (ty_col, pw_col, ly_col, cube_ty_col, cube_ly_col)
    # All raw unit columns — percentages computed after groupby sum (correct weighted avg)
    _COLS = {
        "store":    ("store_oh_units",     "pw_store",         "ly_store",         "store_oh_cube",  "ly_store_cube"),
        "on_order": ("units_ordered",      "pw_units_ordered", "ly_units_ordered", "cube_ordered",   "ly_cube_ordered"),
        "in_dc":    ("in_dc_units",        "pw_dc",            "ly_dc",            "in_dc_cube",     "ly_dc_cube"),
        "it_total": ("it_total_units",     "it_total_pw",      "it_total_ly",      "it_total_cube",  None),
        "it_dc":    ("it_dc_units",        "pw_it_dc",         "ly_it_dc",         "it_dc_cube",     "ly_it_dc_cube"),
        "it_store": ("it_store_units",     "pw_it_store",      "ly_it_store",      "it_store_cube",  "ly_it_store_cube"),
        "backroom": ("backroom_units",     "pw_backroom",      "ly_backroom",      "backroom_cube",  "ly_backroom_cube"),
        "fc":       ("fc_oh_units",        "pw_fc",            "ly_fc",            "fc_oh_cube",     "ly_fc_cube"),
        "yard":     ("on_yard_units",      "pw_yard",          "ly_yard",          "on_yard_cube",   "ly_yard_cube"),
        "total":    ("total_network_units","pw_total",         "ly_total",         None,             None),
    }
    uc, pwc, lyc, cc, lycc = _COLS.get(node, _COLS["store"])
    src = oo if node == "on_order" else inv

    # IT Total: derive combined prior/LY columns from components
    if node == "it_total":
        src = src.copy()
        src["it_total_pw"] = src.get("pw_it_dc",    pd.Series(0.0, index=src.index)).fillna(0) + \
                             src.get("pw_it_store",  pd.Series(0.0, index=src.index)).fillna(0)
        src["it_total_ly"] = src.get("ly_it_dc",    pd.Series(0.0, index=src.index)).fillna(0) + \
                             src.get("ly_it_store",  pd.Series(0.0, index=src.index)).fillna(0)

    if uc not in src.columns:
        return html.P("No data for selected node.", style={"color": "#999", "fontStyle": "italic"})

    # Aggregate only raw unit columns — sum is correct for weighted totals
    sum_cols = [c for c in [uc, pwc, lyc, cc, lycc] if c and c in src.columns]
    grp = src.groupby(["sbu", "OMNI_DEPT_DESC"])[sum_cols].sum().reset_index()

    def _pct_safe(ty, prior):
        """(TY - prior) / prior — returns None when no prior data."""
        if prior is None or prior == 0 or pd.isna(prior): return None
        return (ty - prior) / prior

    def _pc(v):
        if v is None or (isinstance(v, float) and pd.isna(v)): return html.Td("—", style={"textAlign": "right", "color": "#bbb"})
        c = WM_GREEN if v >= 0 else WM_RED
        s = "+" if v >= 0 else ""
        return html.Td(f"{s}{v*100:.1f}%", style={"color": c, "fontWeight": "600", "textAlign": "right"})

    TH = {"background": "#f8f8f8", "fontWeight": "600", "padding": "8px 11px",
          "borderBottom": "2px solid #e0e0e0", "textAlign": "right", "whiteSpace": "nowrap"}
    TD = {"padding": "6px 11px", "borderBottom": "1px solid #f2f2f2"}

    headers = ["SBU", "Department", "Units", "WoW %", "YoY %"]
    if cc: headers += ["Cube (ft³)", "Cube YoY %"]

    rows = [html.Tr([html.Th(h, style={**TH, "textAlign": "left" if i < 2 else "right"})
                     for i, h in enumerate(headers)])]

    prev_sbu = None
    for _, r in grp.sort_values(["sbu", "OMNI_DEPT_DESC"]).iterrows():
        s = r["sbu"]
        if s != prev_sbu:
            sdf = grp[grp["sbu"] == s]
            su  = float(sdf[uc].sum())
            spw = float(sdf[pwc].sum())  if pwc  and pwc  in sdf.columns else 0
            sly = float(sdf[lyc].sum())  if lyc  and lyc  in sdf.columns else 0
            sc  = float(sdf[cc].sum())   if cc   and cc   in sdf.columns else None
            scly= float(sdf[lycc].sum()) if lycc and lycc in sdf.columns else None
            row_cells = [
                html.Td(s, colSpan=2, style={**TD, "fontWeight": "800", "background": "#eef4ff"}),
                html.Td(_fmt(su), style={**TD, "fontWeight": "800", "background": "#eef4ff", "textAlign": "right"}),
                _pc(_pct_safe(su, spw)), _pc(_pct_safe(su, sly)),
            ]
            if cc:
                row_cells += [
                    html.Td(_fmt(sc) if sc else "—", style={**TD, "background": "#eef4ff", "textAlign": "right"}),
                    _pc(_pct_safe(sc, scly)),
                ]
            rows.append(html.Tr(row_cells))
            prev_sbu = s

        uv  = float(r[uc])
        pwv = float(r[pwc])  if pwc  and pwc  in r else 0
        lyv = float(r[lyc])  if lyc  and lyc  in r else 0
        cells = [
            html.Td("", style={**TD}),
            html.Td(r["OMNI_DEPT_DESC"], style={**TD, "paddingLeft": "18px", "color": "#555"}),
            html.Td(_fmt(uv), style={**TD, "textAlign": "right", "fontWeight": "500"}),
            _pc(_pct_safe(uv, pwv)), _pc(_pct_safe(uv, lyv)),
        ]
        if cc:
            cv  = float(r[cc])   if cc   and cc   in r and not pd.isna(r[cc])   else None
            clyv= float(r[lycc]) if lycc and lycc in r and not pd.isna(r[lycc]) else None
            cells += [
                html.Td(_fmt(cv) if cv else "—", style={**TD, "textAlign": "right"}),
                _pc(_pct_safe(cv, clyv)),
            ]
        rows.append(html.Tr(cells))

    return html.Div(
        html.Table(rows, style={"width": "100%", "borderCollapse": "collapse", "fontSize": "0.82rem"}),
        style={"overflowX": "auto"},
    )


@app.callback(
    Output("trend-chart", "figure"),
    Input("trend-node",  "value"),
    Input("trend-range", "value"),
    Input("sbu-filter",  "value"),
    Input("dept-filter", "value"),
)
def update_trend(node, n_wks, sbu, dept):
    n      = int(n_wks)
    is_oo  = (node == "on_order")
    col    = "units_ordered" if is_oo else node
    wk_col = "wm_week" if is_oo else "WM_YR_WK_NBR"
    cur_wk = _cache["cur_oo_wk"] if is_oo else _cache["cur_inv_wk"]

    # Fetch trend data from BQ (or sample fallback)
    if is_oo:
        src = fetch_oo_trend(cur_wk, n_weeks=n)
    else:
        src = fetch_inv_trend(cur_wk, n_weeks=n)

    if sbu  != "All SBUs":        src = src[src["sbu"] == sbu]   if "sbu"  in src.columns else src
    if dept != "All Departments": src = src[src["OMNI_DEPT_DESC"] == dept] if "OMNI_DEPT_DESC" in src.columns else src

    if col not in src.columns or src.empty:
        return go.Figure()

    ty = src.groupby(wk_col)[col].sum().reset_index()
    ty.columns = ["week", "ty"]

    # LY: fetch same N weeks from prior year (offset -100)
    if is_oo:
        ly_src = fetch_oo_trend(cur_wk - 100, n_weeks=n)
    else:
        ly_src = fetch_inv_trend(cur_wk - 100, n_weeks=n)

    if sbu  != "All SBUs":        ly_src = ly_src[ly_src["sbu"] == sbu]   if "sbu"  in ly_src.columns else ly_src
    if dept != "All Departments": ly_src = ly_src[ly_src["OMNI_DEPT_DESC"] == dept] if "OMNI_DEPT_DESC" in ly_src.columns else ly_src

    ly = pd.DataFrame()
    if col in ly_src.columns and not ly_src.empty:
        ly = ly_src.groupby(wk_col)[col].sum().reset_index()
        ly[wk_col] = ly[wk_col] + 100      # shift to TY week numbers for x-axis alignment
        ly.columns = ["week", "ly"]

    merged = ty.merge(ly, on="week", how="left") if not ly.empty else ty

    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=merged["week"].astype(str), y=merged["ty"],
        name="TY", line=dict(color=WM_BLUE, width=2.5),
        mode="lines+markers", marker=dict(size=5),
    ))
    if "ly" in merged.columns:
        fig.add_trace(go.Scatter(
            x=merged["week"].astype(str), y=merged["ly"],
            name="LY (−100 wks)", line=dict(color="#bbb", width=1.5, dash="dot"),
            mode="lines+markers", marker=dict(size=4),
        ))
    fig.update_layout(
        height=240, margin=dict(l=0, r=0, t=6, b=0),
        legend=dict(orientation="h", yanchor="top", y=-0.22, font=dict(size=10)),
        xaxis=dict(title="WM Week", tickfont=dict(size=9), tickangle=-30),
        yaxis=dict(tickfont=dict(size=9)),
        plot_bgcolor="white", paper_bgcolor="white",
        hovermode="x unified",
    )
    fig.update_xaxes(showgrid=True, gridcolor="#f0f0f0")
    fig.update_yaxes(showgrid=True, gridcolor="#f0f0f0")
    return fig


@app.callback(
    Output("insights-output", "children"),
    Output("insights-store",  "data"),
    Input("btn-generate", "n_clicks"),
    State("sbu-filter",  "value"),
    State("dept-filter", "value"),
    prevent_initial_call=True,
)
def generate_insights(n, sbu, dept):
    if not n:
        return dash.no_update, dash.no_update

    from ai.insights import generate_weekly_insights

    inv_cur, oo_cur = _get_frames(sbu, dept)
    cw = _cache["cur_oo_wk"]
    text = generate_weekly_insights(
        inv_week_df=inv_cur, oo_week_df=oo_cur,
        current_wm_week=cw, prior_wm_week=cw - 1,
        current_date=date.today(),
    )
    return html.Pre(text, style={"fontSize": "0.87rem", "lineHeight": "1.75",
                                  "whiteSpace": "pre-wrap", "margin": "0"}), text


@app.callback(
    Output("clipboard", "content"),
    Input("btn-copy",       "n_clicks"),
    State("insights-store", "data"),
    prevent_initial_call=True,
)
def copy_insights(n, text):
    return text or ""


# ── Download Trade Slides PPTX ───────────────────────────────────────────────
_BASE_DIR = os.path.dirname(os.path.abspath(__file__))

def _find_pptx_template():
    """Find pptx_template.pptx first; fallback to most recent readable WK*.pptx.
    Returns (path, bytes) or (None, None)."""
    candidates = [
        os.path.join(_BASE_DIR, "pptx_template.pptx"),
        *sorted(
            glob.glob(os.path.join(_BASE_DIR, "Trade Slides - Inventory WK*.pptx")),
            key=os.path.getmtime, reverse=True,
        ),
    ]
    for path in candidates:
        if not os.path.exists(path):
            continue
        try:
            with open(path, "rb") as f:
                data = f.read()
            if len(data) > 10_000:
                return path, data
        except (PermissionError, OSError):
            continue
    return None, None


def _update_slide1(prs, inv: pd.DataFrame, oo: pd.DataFrame, cur_oo_wk: int):
    """Update all metric text boxes on slide 1 by shape name — no old-value dependency.
    Works every week regardless of what numbers are currently in the template."""
    from datetime import timedelta

    def _n(col): return float(inv[col].sum()) if col in inv.columns else 0.0
    def _o(col): return float(oo[col].sum())  if col in oo.columns  else 0.0
    def _p(a, b): return ((a - b) / b * 100) if b else 0.0

    ty_store = _n("store_oh_units"); ly_store = _n("ly_store"); pw_store = _n("pw_store")
    ty_br    = _n("backroom_units"); ly_br    = _n("ly_backroom"); pw_br  = _n("pw_backroom")
    ty_dc    = _n("in_dc_units");    ly_dc    = _n("ly_dc");    pw_dc    = _n("pw_dc")
    ty_its   = _n("it_store_units"); ly_its   = _n("ly_it_store"); pw_its = _n("pw_it_store")
    ty_itdc  = _n("it_dc_units");    ly_itdc  = _n("ly_it_dc"); pw_itdc  = _n("pw_it_dc")
    ty_fc    = _n("fc_oh_units");    ly_fc    = _n("ly_fc");    pw_fc    = _n("pw_fc")
    ty_yard  = _n("on_yard_units");  ly_yard  = _n("ly_yard");  pw_yard  = _n("pw_yard")
    oo_ty    = _o("units_ordered");  oo_ly    = _o("ly_units_ordered"); oo_pw = _o("pw_units_ordered")
    l13w_ty  = _o("l13w_avg_units"); l13w_ly  = _o("l13w_avg_units_ly")
    ittot_ty = ty_its + ty_itdc;     ittot_ly = ly_its + ly_itdc; ittot_pw = pw_its + pw_itdc
    sf_ty    = ty_store - ty_br;     sf_ly    = ly_store - ly_br; sf_pw  = pw_store - pw_br
    total_ty = _n("total_network_units") + oo_ty
    total_ly = _n("ly_total") + oo_ly
    # Cube values
    oo_cube_ty   = _o("cube_ordered");      oo_cube_ly  = _o("ly_cube_ordered")
    yard_cube_ty = _n("on_yard_cube");      yard_cube_ly = _n("ly_yard_cube")
    dc_cube_ty   = _n("in_dc_cube");        dc_cube_ly  = _n("ly_dc_cube")
    its_cube_ty  = _n("it_store_cube");     its_cube_ly = _n("ly_it_store_cube")
    itdc_cube_ty = _n("it_dc_cube");        itdc_cube_ly= _n("ly_it_dc_cube")
    ittot_cube_ty= its_cube_ty + itdc_cube_ty; ittot_cube_ly = its_cube_ly + itdc_cube_ly
    st_cube_ty   = _n("store_oh_cube");     st_cube_ly  = _n("ly_store_cube")
    br_cube_ty   = _n("backroom_cube");     br_cube_ly  = _n("ly_backroom_cube")
    sf_cube_ty   = st_cube_ty - br_cube_ty; sf_cube_ly = st_cube_ly - br_cube_ly
    fc_cube_ty   = _n("fc_oh_cube");        fc_cube_ly  = _n("ly_fc_cube")

    data_wk  = cur_oo_wk % 100
    trade_wk = data_wk + 1
    meet_date = (date.today() + timedelta(days=7)).strftime("%B %d, %Y").replace(" 0", " ")

    def sm(v):  m = round(v/1e6); return f"{m:,} M" if m >= 1000 else f"{m} M"
    def sm0(v): m = round(v/1e6); return f"{m:,}M" if m >= 1000 else f"{m}M"
    def sl(v):  return (f"+{v:.1f}% vs LY"   if v >= 0 else f"({abs(v):.1f}%) vs LY")
    def sw(v):  return (f"+{v:.1f}% vs LW"   if v >= 0 else f"({abs(v):.1f}%) vs LW")
    def sl_sp(v): return (f"+{v:.1f} % vs LY" if v >= 0 else f"({abs(v):.1f} %) vs LY")
    def sc(v, pct):
        """Cube combined label: 'XXX M ft³ , +X.X% vs. LY'  (K for small values)"""
        sign = "+" if pct >= 0 else ""
        a = abs(v)
        if a >= 1_000_000:
            vfmt = f"{round(v/1e6)} M"
        elif a >= 1_000:
            vfmt = f"{v/1e3:.1f} K"
        else:
            vfmt = f"{int(v)}"
        return f"{vfmt} ft³ , {sign}{pct:.1f}% vs. LY"

    banner = (
        f"6 buckets   ·   TY {total_ty/1e9:.2f} B vs LY {total_ly/1e9:.2f} B   ·   "
        f"{_p(total_ty,total_ly):.1f}% , {(total_ty-total_ly)/1e9:.2f} B YoY   ·   "
        f"{ty_store/1e9:.2f} B store inv ·  {_p(ty_store,ly_store):.1f}%, "
        f"{(ty_store-ly_store)/1e6:.1f} M YoY"
    )
    l13w_txt = (f"L13W Avg: {round(l13w_ty/1e6):,}M/wk "
                f"({'+' if _p(l13w_ty,l13w_ly)>=0 else ''}{_p(l13w_ty,l13w_ly):.1f}% YoY)  ")

    ins1_head = "• On-Order turns positive; BTX floor set complete."
    ins1_body = (f"On-Order {sm0(oo_ty)} ({sl(_p(oo_ty,oo_ly))}) — pipeline above LY. "
                 f"DC releasing at {sm(ty_dc)} ({sl(_p(ty_dc,ly_dc))}). "
                 f"IT→Store {sm(ty_its)} ({sl(_p(ty_its,ly_its))}) — normalizing after BTS surge.")
    ins2_head = f"• {sm(ty_br)} in backroom ({sl(_p(ty_br,ly_br))}) — sustained pull opportunity."
    ins2_body = (f"Salesfloor {sm(sf_ty)} ({sl(_p(sf_ty,sf_ly))}) while backroom builds. "
                 f"CONSUMABLES, PANTRY, and CAC have product staged — "
                 f"pull to shelf without additional receipts. Store OH {sm(ty_store)} ({sl(_p(ty_store,ly_store))}).")

    # Map: shape name → list of paragraph texts (one per paragraph in that shape)
    SHAPE_MAP = {
        "TextBox 98":    [f"Trade Meeting · {meet_date}"],
        "TextBox 2":     [f"Trade Meeting · {meet_date}"],   # slide 2 duplicate
        "TextBox 9":     [banner],
        # On-Order
        "TextBox 19":    [sm0(oo_ty)],
        "TextBox 20":    [sl(_p(oo_ty, oo_ly)), sw(_p(oo_ty, oo_pw))],
        "TextBox 25":    [l13w_txt],
        # On Yard
        "TextBox 35":    [f"{ty_yard/1e6:.1f} M"],
        "TextBox 36":    [sl(_p(ty_yard, ly_yard))],
        "TextBox 21":    [sw(_p(ty_yard, pw_yard))],
        # In DC
        "TextBox 51":    [sm(ty_dc)],
        "TextBox 52":    [sl(_p(ty_dc, ly_dc))],
        "TextBox 24":    [sw(_p(ty_dc, pw_dc))],
        # FC
        "TextBox 117":   [sm(ty_fc)],
        "TextBox 118":   [sl(_p(ty_fc, ly_fc))],
        "TextBox 40":    [sw(_p(ty_fc, pw_fc))],
        # Store OH
        "TextBox 178":   [sm(ty_store)],
        "TextBox 184":   [sl(_p(ty_store, ly_store))],
        "TextBox 185":   [sw(_p(ty_store, pw_store))],
        # Backroom
        "TextBox 188":   [sm(ty_br)],
        "TextBox 189":   [sl_sp(_p(ty_br, ly_br))],
        "TextBox 190":   [sw(_p(ty_br, pw_br))],
        # Salesfloor
        "TextBox 1028":  [sm(sf_ty)],
        "TextBox 1030":  [sl(_p(sf_ty, sf_ly))],
        "TextBox 1032":  [sw(_p(sf_ty, sf_pw))],
        # IT Total
        "TextBox 1034":  [sm(ittot_ty)],
        "TextBox 1036":  [sl(_p(ittot_ty, ittot_ly))],
        "TextBox 1037":  [sw(_p(ittot_ty, ittot_pw))],
        # IT→DC
        "TextBox 1039":  [sm(ty_itdc)],
        "TextBox 1040":  [sl_sp(_p(ty_itdc, ly_itdc))],
        "TextBox 1041":  [sw(_p(ty_itdc, pw_itdc))],
        # IT→Store
        "TextBox 1042":  [sm(ty_its)],
        "TextBox 1043":  [sl_sp(_p(ty_its, ly_its))],
        "TextBox 1044":  [sw(_p(ty_its, pw_its))],
        # ── Cube fields (discovered from template shape inspection) ──────────
        "TextBox 183":   [sc(oo_cube_ty,   _p(oo_cube_ty,   oo_cube_ly))],    # OO cube
        "TextBox 181":   [sc(yard_cube_ty, _p(yard_cube_ty, yard_cube_ly))],  # Yard cube
        "TextBox 180":   [sc(dc_cube_ty,   _p(dc_cube_ty,   dc_cube_ly))],    # DC cube
        "TextBox 1038":  [sc(ittot_cube_ty,_p(ittot_cube_ty,ittot_cube_ly))], # IT Total cube
        "TextBox 54":    [sc(itdc_cube_ty, _p(itdc_cube_ty, itdc_cube_ly))],  # IT→DC cube
        "TextBox 34":    [sc(its_cube_ty,  _p(its_cube_ty,  its_cube_ly))],   # IT→Store cube
        "TextBox 179":   [sc(st_cube_ty,   _p(st_cube_ty,   st_cube_ly))],    # Store OH cube
        "TextBox 41":    [sc(br_cube_ty,   _p(br_cube_ty,   br_cube_ly))],    # Backroom cube
        "TextBox 50":    [sc(sf_cube_ty,   _p(sf_cube_ty,   sf_cube_ly))],    # Salesfloor cube
        "TextBox 3":     [sc(fc_cube_ty,   _p(fc_cube_ty,   fc_cube_ly))],    # FC cube
    }

    # Insights shape (two bullets, each with bold headline + regular detail)
    INSIGHTS_MAP = {
        "Rounded Rectangle 10": [(ins1_head, ins1_body), (ins2_head, ins2_body)],
    }

    changes = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Standard metric boxes
            if shape.name in SHAPE_MAP and shape.has_text_frame:
                new_texts = SHAPE_MAP[shape.name]
                for i, new_text in enumerate(new_texts):
                    if i >= len(shape.text_frame.paragraphs):
                        break
                    para = shape.text_frame.paragraphs[i]
                    if para.runs:
                        para.runs[0].text = new_text
                        for r in para.runs[1:]: r.text = ""
                        changes += 1

            # Insights bullets (bold headline + regular body in same paragraph)
            elif shape.name in INSIGHTS_MAP and shape.has_text_frame:
                bullets = INSIGHTS_MAP[shape.name]
                paras = shape.text_frame.paragraphs
                para_idx = 0
                for (head, body) in bullets:
                    if para_idx >= len(paras): break
                    para = paras[para_idx]
                    runs = para.runs
                    if len(runs) >= 2:
                        runs[0].text = head
                        runs[1].text = body
                        for r in runs[2:]: r.text = ""
                    elif len(runs) == 1:
                        runs[0].text = head + body
                    changes += 1
                    para_idx += 1

    print(f"[PPTX] {changes} shape(s) updated on slide 1")
    return prs


def _update_slide2_chart(prs):
    """Replace Build/Burn chart image on slide 2 with a freshly generated PNG.
    Generates buildburn_chart.png from BQ data via Chrome headless,
    then swaps the existing Picture 7 shape in place."""
    chart_png = os.path.join(_BASE_DIR, "buildburn_chart.png")

    # Generate fresh chart (runs capture_buildburn.py logic inline)
    try:
        sys.path.insert(0, _BASE_DIR)
        import capture_buildburn
        ok = capture_buildburn.capture(out_png=chart_png)
        if not ok:
            print("[PPTX] Build/Burn chart capture failed — skipping slide 2 update")
            return prs
    except Exception as e:
        print(f"[PPTX] Build/Burn chart error: {e} — skipping slide 2 update")
        return prs

    if not os.path.exists(chart_png):
        print("[PPTX] buildburn_chart.png not found — skipping slide 2 update")
        return prs

    try:
        slide2 = prs.slides[1]
        # Find and replace Picture 7
        for shape in list(slide2.shapes):
            if shape.name == "Picture 7":
                left, top, w, h = shape.left, shape.top, shape.width, shape.height
                shape._element.getparent().remove(shape._element)
                slide2.shapes.add_picture(chart_png, left, top, w, h)
                print(f"[PPTX] Slide 2 Build/Burn chart updated ({os.path.getsize(chart_png):,} bytes)")
                break
        else:
            # No Picture 7 found — add chart in default chart area position
            from pptx.util import Inches
            slide2.shapes.add_picture(chart_png, Inches(0), Inches(2), Inches(11), Inches(4))
            print("[PPTX] Slide 2: added new Build/Burn chart (Picture 7 not found)")
    except Exception as e:
        print(f"[PPTX] Slide 2 chart insert error: {e}")

    return prs


def _pptx_replacements(inv: pd.DataFrame, oo: pd.DataFrame, cur_oo_wk: int):
    """Build the same REPLACEMENTS dict as update_pptx.py using cached DataFrames."""
    def _n(col): return float(inv[col].sum()) if col in inv.columns else 0.0
    def _o(col): return float(oo[col].sum())  if col in oo.columns  else 0.0
    def _p(a, b): return ((a - b) / b * 100) if b else 0.0

    # Node values from cache
    ty_store    = _n("store_oh_units");   ly_store    = _n("ly_store");    pw_store    = _n("pw_store")
    ty_br       = _n("backroom_units");   ly_br       = _n("ly_backroom"); pw_br       = _n("pw_backroom")
    ty_dc       = _n("in_dc_units");      ly_dc       = _n("ly_dc");       pw_dc       = _n("pw_dc")
    ty_its      = _n("it_store_units");   ly_its      = _n("ly_it_store"); pw_its      = _n("pw_it_store")
    ty_itdc     = _n("it_dc_units");      ly_itdc     = _n("ly_it_dc");    pw_itdc     = _n("pw_it_dc")
    ty_fc       = _n("fc_oh_units");      ly_fc       = _n("ly_fc");       pw_fc       = _n("pw_fc")
    ty_yard     = _n("on_yard_units");    ly_yard     = _n("ly_yard");     pw_yard     = _n("pw_yard")
    oo_ty       = _o("units_ordered");    oo_ly       = _o("ly_units_ordered"); oo_pw = _o("pw_units_ordered")
    l13w_ty     = _o("l13w_avg_units");   l13w_ly     = _o("l13w_avg_units_ly")
    ty_ittot    = ty_its + ty_itdc;       ly_ittot    = ly_its + ly_itdc;  pw_ittot = pw_its + pw_itdc
    sf_ty       = ty_store - ty_br;       sf_ly       = ly_store - ly_br;  sf_pw = pw_store - pw_br
    total_ty    = _n("total_network_units") + oo_ty
    total_ly    = _n("ly_total") + oo_ly

    from datetime import timedelta
    data_wk  = cur_oo_wk % 100
    trade_wk = data_wk + 1
    # Next Sunday/Monday — simple: just format today+7
    meet_date = (date.today() + timedelta(days=7)).strftime("%B %d, %Y").replace(" 0", " ")

    def sm(v):
        m = round(v / 1e6); return (f"{m:,} M" if m >= 1000 else f"{m} M")
    def sm0(v):
        m = round(v / 1e6); return (f"{m:,}M" if m >= 1000 else f"{m}M")
    def sp(v, s=" vs LY"):
        return (f"+{v:.1f}%{s}" if v >= 0 else f"({abs(v):.1f}%){s}")
    def sl(v): return sp(v, " vs LY")
    def sw(v): return sp(v, " vs LW")

    return {
        # Date
        "Trade Meeting · July 6, 2026":   f"Trade Meeting · {meet_date}",
        # Banner
        "TY 8.32 B vs LY 8.54 B":        f"TY {total_ty/1e9:.2f} B vs LY {total_ly/1e9:.2f} B",
        "-2.5% , -0.22 B YoY":           f"{_p(total_ty,total_ly):.1f}% , {(total_ty-total_ly)/1e9:.2f} B YoY",
        "4.65 B store inv":               f"{ty_store/1e9:.2f} B store inv",
        "-1.4%, -67.1 M YoY":            f"{_p(ty_store,ly_store):.1f}%, {(ty_store-ly_store)/1e6:.1f} M YoY",
        # On-Order
        "1,480M": sm0(oo_ty), "(2.4%) vs LY": sl(_p(oo_ty,oo_ly)), "(6.0%) vs LW": sw(_p(oo_ty,oo_pw)),
        "L13W Avg: 1,370M/wk (+1.6% YoY)": f"L13W Avg: {round(l13w_ty/1e6):,}M/wk ({'+' if _p(l13w_ty,l13w_ly)>=0 else ''}{_p(l13w_ty,l13w_ly):.1f}% YoY)",
        # On Yard
        "0.3 M": f"{ty_yard/1e6:.1f} M", "(94.6%) vs LY": sl(_p(ty_yard,ly_yard)), "(93.7%) vs LW": sw(_p(ty_yard,pw_yard)),
        # In DC
        "1,962 M": sm(ty_dc), "(6.4%) vs LY": sl(_p(ty_dc,ly_dc)), "(6.6%) vs LW": sw(_p(ty_dc,pw_dc)),
        # IT→DC
        "15 M": sm(ty_itdc), "+72.0 % vs LY": sl(_p(ty_itdc,ly_itdc)), "+11.2% vs LW": sw(_p(ty_itdc,pw_itdc)),
        # IT Total
        "167 M": sm(ty_ittot), "+15.7% vs LY": sl(_p(ty_ittot,ly_ittot)), "(0.8%) LW": sw(_p(ty_ittot,pw_ittot)),
        # IT→Store
        "152 M": sm(ty_its), "+12.1 % vs LY": sl(_p(ty_its,ly_its)), "(1.8%) vs LW": sw(_p(ty_its,pw_its)),
        # Store OH
        "4,649 M": sm(ty_store), "(1.4%) vs LY": sl(_p(ty_store,ly_store)), "+0.1% vs LW": sw(_p(ty_store,pw_store)),
        # Backroom
        "422 M": sm(ty_br), "+6.7 % vs LY": sl(_p(ty_br,ly_br)), "+3.9% vs LW": sw(_p(ty_br,pw_br)),
        # Salesfloor
        "4,227 M": sm(sf_ty), "(2.2%) vs LY": sl(_p(sf_ty,sf_ly)), "(0.2%) vs LW": sw(_p(sf_ty,sf_pw)),
        # FC
        "+5.9% vs LY": sl(_p(ty_fc,ly_fc)), "(1.6%) vs LW": sw(_p(ty_fc,pw_fc)),
        # Insights bullet 1
        "• BTX WK23 pipeline confirmed ready.":
            "• On-Order turns positive; BTX floor set complete.",
        "DC depleting -6.4% YoY as inventory flows out; In Transit → Store surging +12.1% YoY with all 8 SBUs above last year, landing in stores within 3–5 days. HARDLINES Stationery BTS floor set staging now.":
            (f"On-Order {sm0(oo_ty)} ({sl(_p(oo_ty,oo_ly))}) — pipeline above LY. "
             f"DC releasing at {sm(ty_dc)} ({sl(_p(ty_dc,ly_dc))}). "
             f"IT→Store {sm(ty_its)} ({sl(_p(ty_its,ly_its))}) — normalizing after BTS surge."),
        # Insights bullet 2
        "• 422M units in backroom — execution pull opportunity.":
            f"• {sm(ty_br)} in backroom ({sl(_p(ty_br,ly_br))}) — sustained pull opportunity.",
        "Store salesfloor is -2.2% YoY (variance in LY data), but backroom inventory is +6.7% YoY. CONSUMABLES, PANTRY, and CAC have product staged and ready to pull to shelf ":
            (f"Salesfloor {sm(sf_ty)} ({sl(_p(sf_ty,sf_ly))}) while backroom builds. "
             f"CONSUMABLES, PANTRY, and CAC have product staged — pull to shelf without additional receipts. "
             f"Store OH {sm(ty_store)} ({sl(_p(ty_store,ly_store))})."),
    }


@server.route("/download-pptx")
def download_trade_slides_route():
    """Flask route — for local use. On Posit Connect use the dcc.Download callback instead."""
    import traceback
    from flask import Response
    try:
        from pptx import Presentation as _Prs
        cur_wk   = _cache["cur_oo_wk"]
        trade_wk = cur_wk % 100 + 1
        filename = f"Trade Slides - Inventory WK{trade_wk}.pptx"
        pptx_path, pptx_bytes = _find_pptx_template()
        if not pptx_path:
            return Response("Template not found. Save pptx_template.pptx in project folder.", status=404)
        prs = _Prs(io.BytesIO(pptx_bytes))
        prs = _update_slide1(prs, _cache["inv"].copy(), _cache["oo"].copy(), cur_wk)
        buf = io.BytesIO(); prs.save(buf); buf.seek(0); data = buf.read()
        print(f"[PPTX] Flask route serving {filename} ({len(data):,} bytes)", flush=True)
        return Response(data,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'})
    except Exception:
        err = traceback.format_exc()
        print(f"[PPTX] Flask route error:\n{err}", flush=True)
        return Response(f"Error: {err}", status=500, mimetype="text/plain")


@app.callback(
    Output("download-pptx", "data"),
    Input("btn-download-pptx", "n_clicks"),
    prevent_initial_call=True,
)
def download_trade_slides_cb(n):
    """Dash callback — works on both local and Posit Connect (uses Dash routing)."""
    import traceback
    if not n:
        return dash.no_update
    try:
        from pptx import Presentation as _Prs
        cur_wk   = _cache["cur_oo_wk"]
        trade_wk = cur_wk % 100 + 1
        filename = f"Trade Slides - Inventory WK{trade_wk}.pptx"
        pptx_path, pptx_bytes = _find_pptx_template()
        if not pptx_path:
            print("[PPTX] No template found — returning no_update", flush=True)
            return dash.no_update
        prs = _Prs(io.BytesIO(pptx_bytes))
        prs = _update_slide1(prs, _cache["inv"].copy(), _cache["oo"].copy(), cur_wk)
        prs = _update_slide2_chart(prs)   # Build/Burn chart on slide 2
        buf = io.BytesIO(); prs.save(buf); buf.seek(0); data = buf.read()
        print(f"[PPTX] Callback serving {filename} ({len(data):,} bytes)", flush=True)
        return dcc.send_bytes(data, filename)
    except Exception:
        print(f"[PPTX] Callback error:\n{traceback.format_exc()}", flush=True)
        return dash.no_update


# ── Run ───────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    app.run(debug=False, host="127.0.0.1", port=8050)
